from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/yanshi")
OUT = ROOT / "PrivEnergyBench/docs/20min_面试汇报_詹绍基_20页图表版.pptx"

NAVY = RGBColor(21, 37, 63)
TEAL = RGBColor(35, 116, 134)
GOLD = RGBColor(191, 145, 65)
TEXT = RGBColor(34, 45, 62)
MUTED = RGBColor(96, 108, 125)
BG = RGBColor(248, 249, 251)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(214, 220, 228)

IMG = {
    "photo": ROOT / "Downloads/詹绍基证件照.jpg",
    "emb_obj": ROOT / "PrivEnergyBench/docs/clean_figures/emb_fixed_objective_clean.png",
    "emb_food": ROOT / "PrivEnergyBench/docs/clean_figures/emb_reward_energy_clean.png",
    "emb_fixed": ROOT / "PrivEnergyBench/docs/clean_figures/emb_fixed_objective_clean.png",
    "emb_heat": ROOT / "PrivEnergyBench/docs/redrawn_figures/emb_food_heatmap_redraw.png",
    "emb_raster": ROOT / "EmbodiedSNNPrototype/outputs/raster.png",
    "ns_pareto": ROOT / "PrivEnergyBench/docs/clean_figures/neuro_pareto_clean.png",
    "ns_3d": ROOT / "PrivEnergyBench/docs/redrawn_figures/neuro_coupling_energy_landscape_redraw.png",
    "ns_2d": ROOT / "PrivEnergyBench/docs/redrawn_figures/neuro_coupling_threshold_energy_lines_redraw.png",
    "med_perf": ROOT / "PrivEnergyBench/docs/clean_figures/med_accuracy_clean.png",
    "med_power": ROOT / "PrivEnergyBench/docs/clean_figures/med_latency_energy_clean.png",
    "med_sparse_mia": ROOT / "PrivEnergyBench/docs/clean_figures/med_mia_auc_clean.png",
    "med_trans": ROOT / "MedSparseSNN/outputs/figures/spiking_transformer_comparison.png",
    "med_cross": ROOT / "MedSparseSNN/outputs/figures/cross_dataset_tradeoff.png",
    "pe_pareto": ROOT / "PrivEnergyBench/outputs/default/pareto.png",
}

SECTIONS = [
    "开场定位",
    "能力证明",
    "对接规划",
    "总结备用",
]


def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0.8)
    return shape


def add_text(slide, text, left, top, width, height, *, size=20, bold=False, color=TEXT,
             font="Noto Sans CJK SC", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color
    return box


def add_bullet_list(slide, items, left, top, width, height, *, size=20, color=TEXT, bullet_color=TEAL, spacing=7):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.name = "Noto Sans CJK SC"
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
        p.bullet = True
        p.level = 0
    for p in tf.paragraphs:
        if p.runs:
            p.runs[0].font.color.rgb = color
    return box


def add_card(slide, title, body, left, top, width, height):
    add_rect(slide, left, top, width, height, CARD, LINE, radius=True)
    add_rect(slide, left, top, width, 0.16, TEAL, TEAL, radius=True)
    add_text(slide, title, left + 0.22, top + 0.24, width - 0.44, 0.42, size=16, bold=True)
    add_text(slide, body, left + 0.22, top + 0.70, width - 0.44, height - 0.92, size=13, color=MUTED)


def add_badge(slide, text, left, top, width=1.45, height=0.34, fill=GOLD, color=RGBColor(255, 255, 255)):
    add_rect(slide, left, top, width, height, fill, fill, radius=True)
    add_text(slide, text, left, top + 0.01, width, height - 0.02, size=11, bold=True, color=color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_result_banner(slide, label, text, left=1.0, top=5.92, width=11.30, height=0.62):
    add_rect(slide, left, top, width, height, NAVY, NAVY, radius=True)
    add_rect(slide, left + 0.10, top + 0.10, 1.10, height - 0.20, GOLD, GOLD, radius=True)
    add_text(slide, label, left + 0.10, top + 0.12, 1.10, height - 0.24, size=10, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, text, left + 1.38, top + 0.10, width - 1.58, height - 0.20, size=14, color=RGBColor(255, 255, 255), valign=MSO_ANCHOR.MIDDLE)


def add_source(slide, text):
    add_text(slide, text, 0.80, 6.80, 10.80, 0.15, size=8, color=MUTED)


def add_key_takeaway(slide, text):
    add_result_banner(slide, "Key Takeaway", text)


def add_picture_contain(slide, path, left, top, width, height):
    if not Path(path).exists():
        add_text(slide, f"Missing: {Path(path).name}", left, top + height / 2 - 0.2, width, 0.4, size=14, align=PP_ALIGN.CENTER)
        return

    with Image.open(path) as img:
        img_width, img_height = img.size

    box_ratio = width / height
    img_ratio = img_width / img_height if img_height else 1

    if img_ratio >= box_ratio:
        draw_width = width
        draw_height = width / img_ratio
        draw_left = left
        draw_top = top + (height - draw_height) / 2
    else:
        draw_height = height
        draw_width = height * img_ratio
        draw_top = top
        draw_left = left + (width - draw_width) / 2

    slide.shapes.add_picture(str(path), Inches(draw_left), Inches(draw_top), width=Inches(draw_width), height=Inches(draw_height))


def add_image_panel(slide, key, left, top, width, height, caption=""):
    add_rect(slide, left, top, width, height, CARD, LINE, radius=True)
    inner_left = left + 0.12
    inner_top = top + 0.12
    inner_width = width - 0.24
    inner_height = height - 0.40
    path = IMG[key]
    if path.exists():
        add_picture_contain(slide, path, inner_left, inner_top, inner_width, inner_height)
    else:
        add_text(slide, f"Missing: {path.name}", inner_left, inner_top + inner_height / 2 - 0.2, inner_width, 0.4, size=14, align=PP_ALIGN.CENTER)
    if caption:
        add_text(slide, caption, left + 0.15, top + height - 0.23, width - 0.30, 0.18, size=10, color=MUTED, align=PP_ALIGN.CENTER)


def add_theme_frame(slide, title, subtitle, page, section_idx):
    add_rect(slide, 0, 0, 13.333, 7.5, BG, BG)
    add_rect(slide, 0, 0, 13.333, 0.74, NAVY, NAVY)
    add_rect(slide, 0.48, 0.74, 0.08, 6.12, TEAL, TEAL)
    add_text(slide, title, 0.75, 0.13, 8.2, 0.30, size=24, bold=True, color=RGBColor(255, 255, 255), font="Noto Serif CJK SC")
    if subtitle:
        add_text(slide, subtitle, 9.2, 0.16, 3.2, 0.24, size=11, color=RGBColor(226, 231, 237), align=PP_ALIGN.RIGHT)
    add_rect(slide, 0.55, 6.97, 12.20, 0.02, LINE, LINE)
    add_text(slide, f"詹绍基 | 面试汇报", 0.75, 7.00, 2.8, 0.20, size=10, color=MUTED)
    add_text(slide, f"{page:02d}", 11.95, 6.98, 0.5, 0.22, size=11, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)

    nav_left = 3.95
    nav_width = 1.85
    for idx, label in enumerate(SECTIONS):
        color = TEAL if idx == section_idx else RGBColor(219, 224, 231)
        text_color = RGBColor(255, 255, 255) if idx == section_idx else MUTED
        add_rect(slide, nav_left + idx * (nav_width + 0.08), 7.00, nav_width, 0.18, color, color, radius=True)
        add_text(slide, label, nav_left + idx * (nav_width + 0.08), 7.00, nav_width, 0.18, size=9, bold=(idx == section_idx), color=text_color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, BG, BG)
    add_rect(slide, 0, 0, 13.333, 1.00, NAVY, NAVY)
    add_text(slide, "面向大规模类脑脉冲神经网络的可信AI研究", 0.90, 0.18, 11.6, 0.52, size=30, bold=True, color=RGBColor(255, 255, 255), font="Noto Serif CJK SC")
    add_text(slide, "科研助理岗位应聘汇报", 0.92, 1.48, 8.0, 0.50, size=28, bold=True, color=NAVY, font="Noto Serif CJK SC")
    add_text(slide, "面试汇报 | 詹绍基", 0.92, 2.08, 4.6, 0.32, size=20, color=MUTED)
    add_text(slide, "长期目标：深耕类脑SNN方向，攻读本方向博士学位", 0.92, 2.36, 8.2, 0.24, size=13, color=MUTED)

    add_rect(slide, 0.92, 2.70, 11.55, 2.25, CARD, LINE, radius=True)
    add_card(slide, "医学×计算机×法学", "跨学科背景\n问题定义更贴近真实约束", 1.10, 2.92, 3.55, 1.72)
    add_card(slide, "SNN / 可信AI全闭环", "从机制验证到评测基准\n完整研究执行链", 4.95, 2.92, 3.55, 1.72)
    add_card(slide, "与实验室高度适配", "方向同频\n可快速启动课题", 8.80, 2.92, 3.55, 1.72)

    add_rect(slide, 0.92, 5.25, 11.55, 0.70, CARD, LINE, radius=True)
    add_text(slide, "20分钟汇报 | 电话：18550954831 | 邮箱：z1309457869@163.com", 1.20, 5.44, 10.9, 0.24, size=16, color=TEXT, align=PP_ALIGN.CENTER)
    add_key_takeaway(slide, "我将读博目标转化为科研助理阶段的高标准执行力：先为团队创造价值，再在同方向深耕成长。")


def content_slide(prs, title, subtitle, page, section_idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_theme_frame(slide, title, subtitle, page, section_idx)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    s = content_slide(prs, "2. 个人核心定位与跨学科优势", "Cross-disciplinary value", 2, 0)
    add_rect(s, 1.00, 1.20, 11.30, 1.20, CARD, LINE, radius=True)
    add_text(s, "苏州大学 医学影像学（辅修计算机） 2019.09-2024.06   |   北京师范大学 法学（二学位） 2024.09-2026.07", 1.20, 1.58, 10.9, 0.32, size=16, align=PP_ALIGN.CENTER)
    add_text(s, "术语说明：脉冲神经网络（Spiking Neural Network, SNN）", 1.20, 2.30, 10.9, 0.20, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_card(s, "计算机算法能力（核心）", "精通Python，熟练掌握SpikingJelly（主力）/SNNtorch/PyTorch\n可独立完成大规模SNN算法设计、多GPU训练与调优\n回答“研究能不能做成”并快速承接科研助理任务", 1.00, 2.75, 3.55, 2.35)
    add_card(s, "医学临床思维（辅助）", "锚定真实场景硬约束\n提升对照实验严谨性\n回答“研究该不该做”", 4.88, 2.75, 3.55, 2.35)
    add_card(s, "法学合规视角（辅助）", "把隐私与问责前置设计\n强化科研诚信与合规意识\n回答“研究能不能用”", 8.76, 2.75, 3.55, 2.35)
    add_text(s, "CET-6 初考595分，可无障碍阅读顶刊顶会文献并参与英文写作，支撑团队论文产出与博士阶段学术输出。", 1.05, 5.45, 11.2, 0.30, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_key_takeaway(s, "以计算机与SNN硬实力为主轴，医学与法学为辅助优势，能直接匹配课题组大规模SNN任务需求。")

    s = content_slide(prs, "3. 核心研究主线与项目闭环", "Large-scale SNN: mechanism -> architecture -> deployment -> benchmark", 3, 1)
    add_rect(s, 1.00, 1.20, 11.30, 1.10, CARD, LINE, radius=True)
    add_text(s, "领域痛点：当前SNN研究多聚焦小体量模型，缺乏面向大规模SNN的训练优化、能效-性能协同与可解释分析框架，同时忽略隐私与真实落地约束。", 1.20, 1.54, 10.9, 0.3, size=16)
    add_card(s, "机制验证", "EmbodiedSNN\n量化可塑性收益-能耗边界", 1.00, 2.70, 2.55, 2.25)
    add_card(s, "架构创新", "Neuro-Symbiosis\n混合架构多目标折中", 3.92, 2.70, 2.55, 2.25)
    add_card(s, "场景落地", "MedSparseSNN\n医疗任务协同评估", 6.84, 2.70, 2.55, 2.25)
    add_card(s, "基础设施", "PrivEnergyBench\n统一评测基准", 9.76, 2.70, 2.55, 2.25)
    add_key_takeaway(s, "四项工作形成大规模SNN可扩展闭环，可直接扩展并对接课题组在训推优化与可信评测上的研究链。")

    s = content_slide(prs, "4. 成果1：具身SNN可塑性机制验证", "EmbodiedSNN", 4, 1)
    add_rect(s, 1.00, 1.22, 3.35, 4.95, CARD, LINE, radius=True)
    add_text(s, "痛点", 1.25, 1.50, 1.0, 0.24, size=16, bold=True)
    add_text(s, "固定权值难以在线适应，\n收益-能耗边界不清。", 1.25, 1.78, 2.9, 0.65, size=15)
    add_text(s, "创新", 1.25, 2.62, 1.0, 0.24, size=16, bold=True)
    add_text(s, "引入奖励调制可塑性\n并量化权衡区间。", 1.25, 2.90, 2.9, 0.65, size=15)
    add_text(s, "证据", 1.25, 3.74, 1.0, 0.24, size=16, bold=True)
    add_text(s, "plastic_readout objective_mean\n优于 structured（-0.4199 vs -0.4839）\n结果以均值±标准差报告。", 1.25, 4.02, 2.95, 1.0, size=13)
    add_image_panel(s, "emb_obj", 4.72, 1.22, 3.70, 4.95, "Fixed vs Plastic Readout Objective Comparison")
    add_image_panel(s, "emb_food", 8.62, 1.22, 3.70, 4.95, "Reward-Energy Trade-off")
    add_source(s, "Source: 本人独立完成的EmbodiedSNNPrototype项目全流程实验输出")
    add_key_takeaway(s, "可塑性显著提升具身智能的在线自适应能力，量化了收益 - 能耗的可优化边界，具备完整的具身 SNN 全流程实践经验，可直接支撑课题组具身 SNN 与大规模模型优化研究。")

    s = content_slide(prs, "5. 成果2：SNN-Transformer混合架构", "Neuro-Symbiosis", 5, 1)
    add_rect(s, 1.00, 1.22, 3.35, 4.95, CARD, LINE, radius=True)
    add_text(s, "痛点", 1.25, 1.50, 1.0, 0.24, size=16, bold=True)
    add_text(s, "单一架构难兼顾\n性能-隐私-能效。", 1.25, 1.78, 2.9, 0.65, size=15)
    add_text(s, "创新", 1.25, 2.62, 1.0, 0.24, size=16, bold=True)
    add_text(s, "提出混合架构，\n验证Pareto折中优势。", 1.25, 2.90, 2.9, 0.65, size=15)
    add_text(s, "证据", 1.25, 3.74, 1.0, 0.24, size=16, bold=True)
    add_text(s, "hybrid: val_acc=0.9612,\nmia_auc=0.4650（均值±标准差另见备用页A）。", 1.25, 4.02, 2.95, 1.0, size=13)
    add_image_panel(s, "ns_pareto", 4.72, 1.22, 3.70, 4.95, "SNN-Transformer Pareto Comparison")
    add_image_panel(s, "ns_3d", 8.62, 1.22, 3.70, 4.95, "Coupling Landscape")
    add_source(s, "Source: 本人独立完成的Neuro-Symbiosis项目全流程实验输出")
    add_key_takeaway(s, "混合架构可实现更优折中，耦合强度是可被搜索优化变量，可支撑课题组大规模SNN训练效率与能效提升。")

    s = content_slide(prs, "6. 成果3：场景验证与评测支撑", "MedSparseSNN + benchmark view", 6, 1)
    add_rect(s, 1.00, 1.22, 3.35, 4.95, CARD, LINE, radius=True)
    add_text(s, "痛点", 1.25, 1.50, 1.0, 0.24, size=16, bold=True)
    add_text(s, "医疗AI需同时满足\n精度、隐私、能效。", 1.25, 1.78, 2.9, 0.65, size=15)
    add_text(s, "创新", 1.25, 2.62, 1.0, 0.24, size=16, bold=True)
    add_text(s, "首次在同一框架下\n协同评估三目标。", 1.25, 2.90, 2.9, 0.65, size=15)
    add_text(s, "结论", 1.25, 3.74, 1.0, 0.24, size=16, bold=True)
    add_text(s, "SNN精度略低但能效显著，\n延迟约0.12 ms/sample，\n各指标均以均值±标准差呈现。", 1.25, 4.02, 2.95, 1.05, size=13)
    add_image_panel(s, "med_perf", 4.72, 1.22, 3.70, 2.35, "Model Accuracy Comparison")
    add_image_panel(s, "med_sparse_mia", 8.62, 1.22, 3.70, 2.35, "MIA AUC Comparison")
    add_image_panel(s, "med_power", 4.72, 3.82, 7.60, 2.35, "Latency-Energy Trade-off")
    add_source(s, "Source: 本人独立完成的MedSparseSNN项目全流程实验输出")
    add_key_takeaway(s, "验证了 SNN 在真实场景的可部署性，为大规模 SNN 在边缘端、行业场景的落地提供了完整的多目标评估范式，可直接服务课题组大规模 SNN 的可信落地研究。")

    s = content_slide(prs, "7. 成果4：统一多目标评测基准", "PrivEnergyBench for large-scale SNN", 7, 1)
    add_rect(s, 1.00, 1.22, 4.20, 4.95, CARD, LINE, radius=True)
    add_text(s, "核心价值", 1.25, 1.52, 1.2, 0.24, size=16, bold=True)
    add_bullet_list(s, [
        "统一性能-隐私-能效指标体系。",
        "多模型同坐标横向比较。",
        "自动统计与Pareto可视化。",
        "可直接支撑大规模SNN架构搜索与多目标优化。",
        "全流程可复现实验报告。",
    ], 1.25, 1.85, 3.6, 1.8, size=14)
    add_text(s, "示例", 1.25, 4.12, 1.2, 0.24, size=16, bold=True)
    add_text(s, "MLP val_acc=0.9208；\nLinear energy=0.0164 mJ/sample。", 1.25, 4.38, 3.7, 0.72, size=14)
    add_image_panel(s, "pe_pareto", 5.55, 1.22, 6.75, 4.95, "Unified Pareto View")
    add_source(s, "Source: 本人独立完成的PrivEnergyBench项目全流程实验输出")
    add_key_takeaway(s, "我已搭建好面向大规模 SNN 的可复用多目标评测底座，可直接承接课题组大规模 SNN 架构搜索、多目标优化与可信 AI 课题，无需从零搭建评测体系。")

    s = content_slide(prs, "8. 项目小结：全闭环科研能力", "Capability proof", 8, 1)
    add_card(s, "问题定义能力", "从真实约束出发而非追逐热点\n可协助PI完成前期调研与问题拆解\n并锻炼博士阶段独立选题能力", 1.00, 1.35, 5.35, 2.15)
    add_card(s, "全链路执行能力", "脚本开发-训练-统计-可视化全流程\n可高效完成PI安排实验任务\n并夯实博士阶段独立实验能力\n具备具身智能场景SNN从环境搭建到评测的完整实践", 6.95, 1.35, 5.35, 2.15)
    add_card(s, "多目标评估能力", "性能-隐私-能效统一分析\n可协助论文结果分析与图表绘制\n并积累学术成果输出能力", 1.00, 3.85, 5.35, 2.15)
    add_card(s, "研究延展能力", "从机制验证到基础设施形成路线\n可独立承接项目子模块\n并培养博士阶段课题规划能力", 6.95, 3.85, 5.35, 2.15)
    add_text(s, "具备具身智能场景 SNN 从环境搭建到评测的完整实践经验（岗位优先录用项）", 1.12, 5.58, 11.0, 0.22, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_key_takeaway(s, "我具备高水平研究全流程基础，并已完成4份实验报告与10+组图表，可直接协助团队论文撰写与投稿。")

    s = content_slide(prs, "9. 核心契合：方法论与实验室同频", "LeCun + Sutton + Lab directions", 9, 2)
    add_card(s, "对齐 LeCun", "世界模型、自监督表征、长时规划\n重抽象预测，不做逐像素重建", 1.00, 1.25, 5.35, 1.90)
    add_card(s, "对齐 Sutton", "通用方法、搜索与学习、可扩展架构\n避免过度手工先验", 6.95, 1.25, 5.35, 1.90)
    add_card(s, "类脑智能方向", "对接 SpiLiFormer (ICCV 2025)\n可塑性机制可用于侧向抑制优化、训练效率提升与能效压缩\n可直接复用实验室现有SpikingJelly代码库，无额外适配成本", 1.00, 3.45, 3.55, 2.45)
    add_card(s, "可信AI与优化", "对接 TEVC 2024 联邦多目标优化\n三目标评测基准可直接接入算法\n实现大规模SNN自动化搜索", 4.88, 3.45, 3.55, 2.45)
    add_card(s, "具身智能", "对接 TEVC 2025 形态进化\n具身SNN平台可复用扩散模型\n实现形态与控制器协同优化", 8.76, 3.45, 3.55, 2.45)
    add_key_takeaway(s, "不是泛泛对接，而是给出可执行的三条落地动作，可直接进入课题组现有研究流程。")

    s = content_slide(prs, "10. 拟承接子课题1：面向大规模SNN的具身智能×世界模型优化", "JEPA-style abstraction for large-scale embodied SNN", 10, 2)
    add_rect(s, 1.00, 1.22, 7.45, 4.95, CARD, LINE, radius=True)
    add_text(s, "目标", 1.25, 1.50, 0.8, 0.24, size=16, bold=True)
    add_text(s, "引入JEPA式抽象预测，提升长时决策稳定性与跨场景泛化。", 2.08, 1.50, 6.1, 0.24, size=15)
    add_text(s, "设计", 1.25, 2.02, 0.8, 0.24, size=16, bold=True)
    add_text(s, "无世界模型 / 像素预测 / 表征预测三组对照；指标：回报、能耗、OOD、迁移。", 2.08, 2.02, 6.1, 0.48, size=15)
    add_text(s, "对接", 1.25, 2.80, 0.8, 0.24, size=16, bold=True)
    add_text(s, "复用现有具身SNN平台，对接实验室具身机器人与进化优化方向。", 2.08, 2.80, 6.1, 0.48, size=15)
    add_text(s, "产出", 1.25, 3.58, 0.8, 0.24, size=16, bold=True)
    add_text(s, "6个月出初步结果；12个月完成主体实验并冲击顶会/顶刊。", 2.08, 3.58, 6.1, 0.24, size=15)

    add_rect(s, 8.75, 1.22, 3.55, 4.95, CARD, LINE, radius=True)
    add_card(s, "6个月", "小环境验证\n+ 基准测试", 8.95, 1.60, 3.15, 1.10)
    add_card(s, "12个月", "核心实验完成\n+ 论文投稿", 8.95, 2.95, 3.15, 1.10)
    add_card(s, "18个月", "多机器人扩展\n+ 方向延展", 8.95, 4.30, 3.15, 1.10)
    add_key_takeaway(s, "该课题可快速启动并直接服务课题组大规模SNN具身方向，填补长时规划能力关键空白。")

    s = content_slide(prs, "11. 拟承接子课题2：大规模SNN在可信约束下的多目标协同优化", "Federated + DP + multi-objective search for large-scale SNN", 11, 2)
    add_rect(s, 1.00, 1.22, 7.45, 4.95, CARD, LINE, radius=True)
    add_text(s, "目标", 1.25, 1.50, 0.8, 0.24, size=16, bold=True)
    add_text(s, "在联邦/分布式场景下统一优化性能、隐私、能效。", 2.08, 1.50, 6.1, 0.24, size=15)
    add_text(s, "创新", 1.25, 2.02, 0.8, 0.24, size=16, bold=True)
    add_text(s, "把DP预算、MIA风险、能耗成本由后置评估前移为优化目标。", 2.08, 2.02, 6.1, 0.48, size=15)
    add_text(s, "对接", 1.25, 2.80, 0.8, 0.24, size=16, bold=True)
    add_text(s, "复用PrivEnergyBench，对接实验室联邦学习、差分隐私、多目标进化优化积累。", 2.08, 2.80, 6.1, 0.48, size=15)
    add_text(s, "产出", 1.25, 3.58, 0.8, 0.24, size=16, bold=True)
    add_text(s, "3个月工具包；9个月核心结果；15个月1-2篇高水平论文与开源发布。", 2.08, 3.58, 6.1, 0.48, size=15)

    add_rect(s, 8.75, 1.22, 3.55, 4.95, CARD, LINE, radius=True)
    add_card(s, "3个月", "联邦场景适配\n+ 首个工具包", 8.95, 1.60, 3.15, 1.10)
    add_card(s, "9个月", "融合多目标搜索\n+ 核心实验", 8.95, 2.95, 3.15, 1.10)
    add_card(s, "15个月", "医疗场景验证\n+ 论文与开源", 8.95, 4.30, 3.15, 1.10)
    add_key_takeaway(s, "该课题与实验室可信AI和多目标优化主线高度重合，具备最快落地路径。")

    s = content_slide(prs, "12. 入职后0-6个月执行规划", "Short-term execution roadmap", 12, 2)
    add_rect(s, 1.00, 1.20, 11.30, 5.25, CARD, LINE, radius=True)
    add_card(s, "第1个月", "完成团队对接、熟悉SNN项目框架\n配合完成基础实验/数据/图表\n同步完成基准与算法初步适配\n对齐SpikingJelly代码规范与项目流程\n并梳理前沿文献明确博士方向", 1.20, 1.62, 2.55, 4.45)
    add_card(s, "第2-3个月", "全力完成联邦场景适配\n输出首个可复用评测工具包\n完成文献梳理、数据整理与论文图表绘制\n协助PI完成论文初稿撰写与润色", 4.05, 1.62, 2.55, 4.45)
    add_card(s, "第4-5个月", "完成课题1核心对照实验\n同步推进课题2算法融合\n独立负责对应子模块全流程\n并整理申博前期材料", 6.90, 1.62, 2.55, 4.45)
    add_card(s, "第6个月", "完成两课题阶段总结\n形成1篇工作论文\n全力配合投稿工作\n在PI指导下明确博士选题", 9.75, 1.62, 2.35, 4.45)
    add_key_takeaway(s, "6个月内可完成平台对接、工具与论文阶段成果，既为团队形成实质贡献，也完成申博核心基础积累。")

    s = content_slide(prs, "13. 中长期研究规划（1-3年）", "Year 1 -> Year 3", 13, 2)
    add_card(s, "第1年（科研助理核心期）", "100%完成PI安排科研任务，协助实验/分析/写作\n产出2-3篇高质量合作论文\n能力达标后独立负责1-2个子模块\n在PI指导下优先推进本实验室博士申请", 1.00, 1.65, 3.55, 4.65)
    add_card(s, "第2年（博士衔接期）", "若成功申请本实验室博士：无缝衔接并持续深耕\n基于科研助理成果完成博士课题完整设计与申报\n若申请其他同方向顶尖实验室博士：\n先完成岗位收尾、项目交接、代码文档沉淀\n确保课题组项目与投稿不受影响", 4.88, 1.65, 3.55, 4.65)
    add_card(s, "第3年（长期深耕期）", "无论博士去向，持续深耕类脑SNN\n与课题组保持长期学术合作\n共同推进成果产出与落地", 8.76, 1.65, 3.55, 4.65)
    add_key_takeaway(s, "我的1-3年规划与课题组发展同向：团队产出始终第一优先，并在此基础上稳步推进博士成长。")

    s = content_slide(prs, "14. 为什么是我：5个不可替代优势", "Final value proposition", 14, 3)
    add_card(s, "方向高度匹配", "研究主线与实验室三大方向同频\n已有成果可无缝衔接", 1.00, 1.35, 5.35, 2.15)
    add_card(s, "全闭环科研能力", "独立完成问题定义到成果输出\n入组即可启动课题", 6.95, 1.35, 5.35, 2.15)
    add_card(s, "明确学术规划与强自驱", "以深耕大规模SNN并攻读博士为长期目标\n以最高标准完成科研助理工作\n实现个人成长与团队产出双赢", 1.00, 3.55, 3.55, 2.15)
    add_card(s, "可落地研究理念", "可复现、可验证、可部署\n与实验室标准一致", 4.88, 3.55, 3.55, 2.15)
    add_card(s, "稀缺跨学科背景", "医学+计算机+法学\n拓展医疗可信AI边界", 8.76, 3.55, 3.55, 2.15)
    add_key_takeaway(s, "我带来的不是零散项目，而是一条可与课题组深度融合、共同成长的学术路线，可胜任科研助理并成长为博士骨干。")

    s = content_slide(prs, "15. 致谢", "Q&A", 15, 3)
    add_rect(s, 1.00, 1.55, 11.30, 4.65, CARD, LINE, radius=True)
    add_text(s, "感谢各位老师的聆听！欢迎批评指正", 1.25, 2.08, 10.8, 0.6, size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Noto Serif CJK SC")
    add_text(s, "詹绍基 | 电话：18550954831 | 邮箱：z1309457869@163.com", 1.25, 3.10, 10.8, 0.3, size=20, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, "可提供全部实验代码、报告与原始数据，确保全流程可复现。", 1.25, 4.00, 10.8, 0.3, size=18, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(s, "长期目标是深耕类脑SNN方向并攻读本方向博士，期待在老师指导下成长并为团队创造长期价值。", 1.25, 4.55, 10.8, 0.25, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_key_takeaway(s, "所有结果均可复现，可即时提供代码与实验明细，也可随时沟通学术规划与岗位安排。")

    s = content_slide(prs, "备用页A：实验详细数据汇总", "Backup for technical questions", 16, 3)
    add_rect(s, 1.00, 1.20, 11.30, 5.30, CARD, LINE, radius=True)
    add_text(s, "用于评委追问时快速调用", 1.30, 1.52, 4.0, 0.24, size=16, color=MUTED)
    add_text(s, "核心技术栈：SpikingJelly（主力）/ SNNtorch / PyTorch / NumPy / Matplotlib", 1.30, 1.78, 10.7, 0.24, size=14, color=TEXT)
    add_bullet_list(s, [
        "EmbodiedSNN：objective / food / energy / spike activity 关键统计。",
        "Neuro-Symbiosis：baseline对照、Pareto前沿与耦合敏感性。",
        "MedSparseSNN：accuracy、MIA AUC、power-latency 全量表。",
        "PrivEnergyBench：多seed mean/std/CI 与可复现命令入口。",
        "大规模SNN技术理解：稀疏训练优化、多GPU并行训推、能效-性能-隐私多目标协同优化。",
        "已完成SpikingJelly多GPU训练环境搭建与测试，可直接承接大规模SNN训推任务。",
        "落地方向：用SpikingJelly训练栈承接大规模SNN实验，直接复用到课题组现有项目流程。",
        "可按“模型-数据-指标-结论-复现路径”五元组现场展开。",
    ], 1.30, 2.05, 10.7, 2.9, size=14)
    add_image_panel(s, "pe_pareto", 1.30, 4.62, 10.7, 1.75, "Aggregated Pareto for quick Q&A")
    add_key_takeaway(s, "备用页用于回答统计显著性、复现路径与短板解释，不参与正式主讲计时。")

    s = content_slide(prs, "备用页B：核心参考文献", "Backup references", 17, 3)
    add_rect(s, 1.00, 1.20, 11.30, 5.30, CARD, LINE, radius=True)
    add_bullet_list(s, [
        "[Lab] SpiLiFormer, ICCV 2025；Morphology Evolution, IEEE TEVC 2025。",
        "[Lab] Federated many-task Bayesian optimization, IEEE TEVC 2024。",
        "[Lab] DP-FSAEA, IEEE TEVC 2024；Post-LLM roadmap, Engineering 2024。",
        "[Method] LeCun, A Path Towards Autonomous Machine Intelligence, 2022。",
        "[Method] Sutton, The Bitter Lesson, 2019。",
        "[Mine] EmbodiedSNNPrototype / Neuro-Symbiosis / MedSparseSNN / PrivEnergyBench 实验报告。",
    ], 1.25, 1.70, 10.8, 4.6, size=14)
    add_key_takeaway(s, "主讲仅引用3个关键对接点，完整文献放在备用页供评委追问时展开。")

    s = content_slide(prs, "备用页C：博士阶段拟研究方向", "Backup for PhD-plan questions", 18, 3)
    add_rect(s, 1.00, 1.20, 11.30, 5.30, CARD, LINE, radius=True)
    add_card(s, "方向1", "大规模脉冲神经网络\n能效优化与稀疏性提升", 1.20, 1.65, 3.45, 2.10)
    add_card(s, "方向2", "世界模型赋能\n具身SNN长时规划", 4.95, 1.65, 3.45, 2.10)
    add_card(s, "方向3", "可信约束下大规模SNN\n多目标协同优化", 8.70, 1.65, 3.45, 2.10)
    add_rect(s, 1.20, 4.05, 10.95, 2.20, CARD, LINE, radius=True)
    add_text(s, "衔接优势", 1.45, 4.36, 1.2, 0.24, size=16, bold=True)
    add_text(s, "以上方向全部基于科研助理阶段工作积累，可无缝衔接、避免重复投入、快速形成博士阶段成果。", 2.70, 4.36, 8.95, 0.50, size=15)
    add_text(s, "预期价值", 1.45, 5.00, 1.2, 0.24, size=16, bold=True)
    add_text(s, "博士阶段研究将反向拓展课题组研究边界，与课题组形成长期协同。", 2.70, 5.00, 8.95, 0.32, size=15)
    add_key_takeaway(s, "博士规划完全对齐课题组三大主线，科研助理阶段投入可直接转化为长期共同产出。")

    prs.save(str(OUT))
    print(f"saved: {OUT}")


if __name__ == "__main__":
    build()
