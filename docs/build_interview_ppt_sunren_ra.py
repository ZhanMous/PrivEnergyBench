from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/home/yanshi")
OUT = ROOT / "PrivEnergyBench/docs/下午面试_西湖大学孙仁组_科研助理版.pptx"

NAVY = RGBColor(18, 37, 64)
TEAL = RGBColor(33, 116, 133)
GOLD = RGBColor(191, 145, 65)
TEXT = RGBColor(34, 45, 62)
MUTED = RGBColor(96, 108, 125)
BG = RGBColor(248, 249, 251)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(214, 220, 228)

SECTIONS = ["实验室理解", "岗位匹配", "执行方案", "收尾问答"]


def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0.8)
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    size=20,
    bold=False,
    color=TEXT,
    font="Noto Sans CJK SC",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color
    return box


def add_bullets(slide, items, left, top, width, height, *, size=15, color=TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Noto Sans CJK SC"
        p.font.color.rgb = color
        p.space_after = Pt(5)
    return box


def add_card(slide, title, body, left, top, width, height):
    add_rect(slide, left, top, width, height, CARD, LINE, radius=True)
    add_rect(slide, left, top, width, 0.14, TEAL, TEAL, radius=True)
    add_text(slide, title, left + 0.20, top + 0.20, width - 0.4, 0.34, size=15, bold=True)
    add_text(slide, body, left + 0.20, top + 0.58, width - 0.4, height - 0.74, size=13, color=MUTED)


def add_takeaway(slide, text):
    add_rect(slide, 1.00, 5.92, 11.30, 0.62, NAVY, NAVY, radius=True)
    add_rect(slide, 1.10, 6.02, 1.15, 0.42, GOLD, GOLD, radius=True)
    add_text(
        slide,
        "Key Takeaway",
        1.10,
        6.04,
        1.15,
        0.38,
        size=10,
        bold=True,
        color=RGBColor(255, 255, 255),
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        text,
        2.38,
        6.02,
        9.70,
        0.42,
        size=13,
        color=RGBColor(255, 255, 255),
        valign=MSO_ANCHOR.MIDDLE,
    )


def frame(slide, title, subtitle, page, sec_idx):
    add_rect(slide, 0, 0, 13.333, 7.5, BG, BG)
    add_rect(slide, 0, 0, 13.333, 0.74, NAVY, NAVY)
    add_rect(slide, 0.48, 0.74, 0.08, 6.12, TEAL, TEAL)
    add_text(slide, title, 0.75, 0.13, 8.6, 0.30, size=24, bold=True, color=RGBColor(255, 255, 255), font="Noto Serif CJK SC")
    if subtitle:
        add_text(slide, subtitle, 8.95, 0.15, 3.45, 0.24, size=11, color=RGBColor(226, 231, 237), align=PP_ALIGN.RIGHT)

    add_rect(slide, 0.55, 6.97, 12.20, 0.02, LINE, LINE)
    add_text(slide, "詹绍基 | 科研助理岗位面试", 0.75, 7.00, 4.1, 0.20, size=10, color=MUTED)
    add_text(slide, f"{page:02d}", 11.95, 6.98, 0.5, 0.22, size=11, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)

    nav_left = 3.95
    nav_width = 1.85
    for idx, label in enumerate(SECTIONS):
        color = TEAL if idx == sec_idx else RGBColor(219, 224, 231)
        tcolor = RGBColor(255, 255, 255) if idx == sec_idx else MUTED
        add_rect(slide, nav_left + idx * (nav_width + 0.08), 7.00, nav_width, 0.18, color, color, radius=True)
        add_text(
            slide,
            label,
            nav_left + idx * (nav_width + 0.08),
            7.00,
            nav_width,
            0.18,
            size=9,
            bold=(idx == sec_idx),
            color=tcolor,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )


def content(prs, title, subtitle, page, sec_idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    frame(s, title, subtitle, page, sec_idx)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 13.333, 7.5, BG, BG)
    add_rect(s, 0, 0, 13.333, 1.00, NAVY, NAVY)
    add_text(
        s,
        "西湖大学孙仁教授课题组科研助理面试汇报",
        0.92,
        0.20,
        11.8,
        0.48,
        size=30,
        bold=True,
        color=RGBColor(255, 255, 255),
        font="Noto Serif CJK SC",
    )
    add_text(s, "岗位导向：病毒学/免疫学方向实验与论文协同执行", 0.95, 1.48, 8.4, 0.38, size=22, bold=True, color=NAVY)
    add_text(s, "詹绍基 | 下午面试专版", 0.95, 2.06, 5.4, 0.30, size=18, color=MUTED)

    add_card(s, "我这次汇报只回答三件事", "我如何理解实验室方向\n我如何匹配岗位职责\n我入组后如何快速产出", 0.95, 2.75, 3.65, 1.75)
    add_card(s, "岗位核心场景", "生物实验执行\n数据记录与统计\n论文与综述协助", 4.92, 2.75, 3.65, 1.75)
    add_card(s, "目标", "下午面试可直接沟通\n入组后的执行计划与分工", 8.89, 2.75, 3.65, 1.75)

    add_rect(s, 0.95, 5.18, 11.58, 0.78, CARD, LINE, radius=True)
    add_text(s, "联系方式：18550954831 | z1309457869@163.com", 1.22, 5.45, 11.1, 0.24, size=16, align=PP_ALIGN.CENTER)
    add_takeaway(s, "这是一份岗位导向汇报：以课题组研究产出为目标，强调执行力与协同价值。")

    # 2 lab overview
    s = content(prs, "2. 对实验室的理解", "Prof. Ren Sun Lab overview", 2, 0)
    add_rect(s, 1.00, 1.20, 11.30, 4.95, CARD, LINE, radius=True)
    add_text(s, "里程碑认识", 1.25, 1.55, 1.8, 0.24, size=17, bold=True)
    add_bullets(
        s,
        [
            "2022年加入西湖大学，长期聚焦病毒学与免疫学，强调系统生物学整合。",
            "实现单核苷酸精度病毒功能图谱解析，并用于减毒疫苗精准设计。",
            "在RNA病毒结构解析和病毒-宿主相互作用机理方面具有高影响力成果。",
        ],
        1.25,
        1.92,
        10.8,
        1.55,
        size=15,
    )
    add_text(s, "我理解的课题组工作方式", 1.25, 3.67, 2.8, 0.24, size=17, bold=True)
    add_bullets(
        s,
        [
            "实验密度高：高通量策略 + 连续迭代验证。",
            "数据标准高：记录可追溯、统计可复查、结论可复现。",
            "论文节奏快：实验和写作同步推进，强调协同效率。",
        ],
        1.25,
        4.02,
        10.8,
        1.60,
        size=15,
    )
    add_takeaway(s, "我会按‘实验可复现、记录可追溯、写作可交付’三条标准进入团队节奏。")

    # 3 directions
    s = content(prs, "3. 三条研究方向拆解", "What the group does next", 3, 0)
    add_card(
        s,
        "方向A：病毒基因组精准功能定位与疫苗",
        "关键任务：高通量突变体分析、功能区定位、候选位点改造\n代表成果：Science 2018, mBio 2017, PNAS 2017",
        1.00,
        1.35,
        11.30,
        1.55,
    )
    add_card(
        s,
        "方向B：感染与接种后免疫反应系统评估",
        "关键任务：抗体特异性评估、T/B细胞响应分析、免疫组学流程构建\n应用延展：肿瘤、自身免疫、神经免疫",
        1.00,
        3.10,
        11.30,
        1.55,
    )
    add_card(
        s,
        "方向C：蛋白结构与功能设计",
        "关键任务：体外进化 + 结构解析 + 质谱，设计和优化蛋白/多肽结合体\n参考文献：Cell 2019, Nature 2018/2017, eLife 2016",
        1.00,
        4.85,
        11.30,
        1.55,
    )
    add_takeaway(s, "三条方向共性是‘实验-数据-机制-产出’闭环，这与科研助理岗位高度一致。")

    # 4 job mapping
    s = content(prs, "4. 岗位职责逐条对齐", "RA responsibilities mapping", 4, 1)
    add_card(
        s,
        "职责1：参与实验",
        "分子克隆、流式细胞、小鼠实验等，需要稳定执行与规范操作。",
        1.00,
        1.35,
        3.55,
        2.25,
    )
    add_card(
        s,
        "职责2：记录与分析",
        "按时提交实验记录与报告，完成数据统计分析并持续学习新技术。",
        4.88,
        1.35,
        3.55,
        2.25,
    )
    add_card(
        s,
        "职责3：写作与协作",
        "协助撰写研究论文/综述，完成PI交办的实验室协同任务。",
        8.76,
        1.35,
        3.55,
        2.25,
    )
    add_rect(s, 1.00, 3.95, 11.30, 2.20, CARD, LINE, radius=True)
    add_text(s, "岗位关键词", 1.28, 4.28, 1.4, 0.24, size=17, bold=True)
    add_bullets(
        s,
        [
            "实验可靠性：操作规范、记录完整、问题可追溯。",
            "沟通协作：与博士后/研究生协同推进，按时间点交付。",
            "英文能力：文献检索、结果整理、英文写作配合。",
        ],
        2.90,
        4.20,
        8.9,
        1.45,
        size=15,
    )
    add_takeaway(s, "我会把岗位要求拆成可执行动作：按周交付实验、记录、分析和写作支持。")

    # 5 fit strengths
    s = content(prs, "5. 我的匹配优势（面向岗位）", "What I can bring quickly", 5, 1)
    add_card(
        s,
        "优势1：执行纪律",
        "临床实习阶段形成标准化记录和流程意识，适应高要求实验室管理。",
        1.00,
        1.35,
        5.35,
        2.15,
    )
    add_card(
        s,
        "优势2：数据与文档能力",
        "可快速完成结果整理、表格汇总、图示规范化与报告输出。",
        6.95,
        1.35,
        5.35,
        2.15,
    )
    add_card(
        s,
        "优势3：英文文献能力",
        "CET-6 595分，能高频检索与阅读文献并支持写作协作。",
        1.00,
        3.85,
        5.35,
        2.15,
    )
    add_card(
        s,
        "优势4：主动学习与配合",
        "可在PI指导下快速补齐特定实验模块，并按节点稳定交付。",
        6.95,
        3.85,
        5.35,
        2.15,
    )
    add_takeaway(s, "我的定位是‘高执行+高配合’科研助理，先把团队任务做扎实，再持续提升实验深度。")

    # 6 risk and commitment
    s = content(prs, "6. 风险识别与承诺", "How I reduce PI concerns", 6, 1)
    add_rect(s, 1.00, 1.20, 11.30, 4.95, CARD, LINE, radius=True)
    add_text(s, "PI最关注的三件事", 1.25, 1.55, 2.2, 0.24, size=17, bold=True)
    add_bullets(
        s,
        [
            "会不会把岗位当跳板？",
            "实验记录是否严谨、是否可复查？",
            "是否能长期稳定协作并按时交付？",
        ],
        1.25,
        1.92,
        4.6,
        1.50,
        size=15,
    )
    add_text(s, "我的承诺", 6.10, 1.55, 1.6, 0.24, size=17, bold=True)
    add_bullets(
        s,
        [
            "团队任务优先：先完成PI安排的实验与写作协同，再谈个人发展。",
            "交付可追溯：实验记录、数据脚本、分析文档按规范沉淀。",
            "沟通可预期：周报+节点复盘，问题提前暴露不过夜。",
        ],
        6.10,
        1.92,
        5.9,
        1.90,
        size=15,
    )
    add_rect(s, 1.25, 4.10, 10.90, 1.65, CARD, LINE, radius=True)
    add_text(s, "一句话", 1.55, 4.38, 1.0, 0.24, size=16, bold=True)
    add_text(s, "我把科研助理岗位视为‘高质量贡献阶段’：用稳定产出换取长期信任与成长。", 2.70, 4.38, 8.9, 0.55, size=15)
    add_takeaway(s, "我会用可追踪交付和长期配合，降低PI对‘执行不稳、协作不稳’的核心顾虑。")

    # 7 plan 0-90d
    s = content(prs, "7. 入组0-90天执行计划", "Immediate execution roadmap", 7, 2)
    add_card(
        s,
        "第1-2周",
        "熟悉课题组SOP与安全规范\n跟随完成实验流程演示\n建立记录模板",
        1.00,
        1.35,
        2.60,
        4.80,
    )
    add_card(
        s,
        "第3-4周",
        "独立完成指定基础实验模块\n按时提交实验记录与数据表\n完成首轮文献清单",
        3.95,
        1.35,
        2.60,
        4.80,
    )
    add_card(
        s,
        "第2个月",
        "承担稳定实验批次执行\n配合统计分析与图表整理\n参与论文材料准备",
        6.90,
        1.35,
        2.60,
        4.80,
    )
    add_card(
        s,
        "第3个月",
        "形成可复用实验/分析文档\n协助初稿撰写或综述整理\n提出下一步优化建议",
        9.85,
        1.35,
        2.35,
        4.80,
    )
    add_takeaway(s, "目标不是‘试试看’，而是90天内完成可见交付并进入稳定产出状态。")

    # 8 writing workflow
    s = content(prs, "8. 论文协同与数据交付流程", "From bench to manuscript", 8, 2)
    add_rect(s, 1.00, 1.20, 11.30, 4.95, CARD, LINE, radius=True)
    add_card(s, "Step 1\n实验记录", "原始记录-版本管理-异常标注", 1.20, 1.65, 2.35, 1.70)
    add_card(s, "Step 2\n数据整理", "统一命名-统计表-可追溯脚本", 3.95, 1.65, 2.35, 1.70)
    add_card(s, "Step 3\n图表输出", "图例/单位/标题标准化", 6.70, 1.65, 2.35, 1.70)
    add_card(s, "Step 4\n写作协同", "结果段落草稿+参考文献整理", 9.45, 1.65, 2.35, 1.70)

    add_text(s, "质量控制", 1.25, 3.85, 1.2, 0.24, size=17, bold=True)
    add_bullets(
        s,
        [
            "每周固定交付一次‘实验-数据-图表’打包文件。",
            "所有结果都可从原始记录追溯到最终图表。",
            "英文材料按模板提交，减少团队返工。",
        ],
        2.65,
        3.82,
        8.9,
        1.35,
        size=15,
    )
    add_takeaway(s, "我可直接承担‘实验执行+数据整理+写作支持’的闭环协作角色。")

    # 9 expected value
    s = content(prs, "9. 我能给团队带来的具体价值", "Team-first value", 9, 2)
    add_card(
        s,
        "短期价值（1-3个月）",
        "补充实验执行产能\n提升记录与分析规范度\n缓解研究主力在重复流程上的压力",
        1.00,
        1.35,
        3.55,
        4.80,
    )
    add_card(
        s,
        "中期价值（3-6个月）",
        "承担稳定子模块执行\n形成可复用流程文档\n加速论文材料整理节奏",
        4.88,
        1.35,
        3.55,
        4.80,
    )
    add_card(
        s,
        "长期价值（6个月+）",
        "形成高配合度协作关系\n持续提升技术深度\n成为可依赖的执行骨干",
        8.76,
        1.35,
        3.55,
        4.80,
    )
    add_takeaway(s, "我的核心价值是‘稳定交付+规范协作+持续进步’，优先服务课题组产出。")

    # 10 qna prep
    s = content(prs, "10. 高频提问预答", "Afternoon interview quick answers", 10, 3)
    add_rect(s, 1.00, 1.20, 11.30, 4.95, CARD, LINE, radius=True)
    add_card(
        s,
        "Q1：你如何保证实验严谨？",
        "A：按SOP执行+全流程记录+可追溯数据管理，结果可复查。",
        1.20,
        1.55,
        5.20,
        1.30,
    )
    add_card(
        s,
        "Q2：能否承担论文协助？",
        "A：可承担文献检索、数据表、图表规范化与结果段草稿协作。",
        6.10,
        1.55,
        5.20,
        1.30,
    )
    add_card(
        s,
        "Q3：你和岗位最匹配的点？",
        "A：执行纪律强、交付稳定、沟通及时，能快速进入团队节奏。",
        1.20,
        3.05,
        5.20,
        1.30,
    )
    add_card(
        s,
        "Q4：你的长期规划会影响岗位吗？",
        "A：不会。团队任务始终第一优先，先交付价值再谈个人成长。",
        6.10,
        3.05,
        5.20,
        1.30,
    )
    add_takeaway(s, "面试回答统一原则：先团队价值、再个人规划；先可执行、再愿景。")

    # 11 thanks
    s = content(prs, "11. 致谢", "Thanks", 11, 3)
    add_rect(s, 1.00, 1.55, 11.30, 4.65, CARD, LINE, radius=True)
    add_text(s, "感谢各位老师聆听，欢迎批评指正", 1.25, 2.08, 10.8, 0.6, size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Noto Serif CJK SC")
    add_text(s, "詹绍基 | 18550954831 | z1309457869@163.com", 1.25, 3.20, 10.8, 0.30, size=20, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, "可在问答阶段展开：实验记录模板、数据统计流程、论文协同样例", 1.25, 4.10, 10.8, 0.25, size=15, color=MUTED, align=PP_ALIGN.CENTER)
    add_takeaway(s, "期待加入孙仁教授课题组，以稳定执行和高质量协同支持团队研究产出。")

    # 12 backup
    s = content(prs, "12. 备用页：文献与方向速览", "Backup", 12, 3)
    add_rect(s, 1.00, 1.20, 11.30, 5.30, CARD, LINE, radius=True)
    add_bullets(
        s,
        [
            "病毒功能图谱与疫苗：Du et al., Science 2018; Du et al., mBio 2017; Qi et al., PNAS 2017。",
            "减毒疫苗策略：Du et al., Science 2018; Gong et al., Cell Host Microbe 2016; NPJ Vaccines 2020。",
            "蛋白结构功能：Dong et al., Cell 2019; Dai et al., Nature 2018/2017; Wu et al., eLife 2016。",
            "面试可追问点：我如何参与实验执行、数据分析和论文协同的具体分工。",
        ],
        1.25,
        1.70,
        10.8,
        4.5,
        size=14,
    )
    add_takeaway(s, "主讲聚焦岗位价值，备用页用于快速响应文献与方向追问。")

    prs.save(str(OUT))
    print(f"saved: {OUT}")


if __name__ == "__main__":
    build()
