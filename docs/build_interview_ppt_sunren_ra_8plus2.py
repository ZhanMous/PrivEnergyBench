from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/home/yanshi")
OUT = ROOT / "PrivEnergyBench/docs/下午面试_孙仁组_RA_8页正文2备用.pptx"

NAVY = RGBColor(18, 37, 64)
TEAL = RGBColor(33, 116, 133)
GOLD = RGBColor(191, 145, 65)
TEXT = RGBColor(34, 45, 62)
MUTED = RGBColor(96, 108, 125)
BG = RGBColor(248, 249, 251)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(214, 220, 228)
SOFT_TEAL_BG = RGBColor(230, 242, 245)

SECTIONS = ["核心匹配", "能力证明", "执行规划", "收尾备用"]


def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0.8)
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    size=24,
    bold=False,
    color=TEXT,
    font="Noto Sans CJK SC",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = font
    p.font.color.rgb = color
    return box


def add_bullets(slide, items, left, top, width, height, *, size=20, color=TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.bullet = True
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Noto Sans CJK SC"
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return box


def add_card(slide, title, body, left, top, width, height):
    add_rect(slide, left, top, width, height, CARD, LINE, radius=True)
    add_rect(slide, left, top, width, 0.14, TEAL, TEAL, radius=True)
    add_text(slide, title, left + 0.20, top + 0.20, width - 0.4, 0.32, size=18, bold=True)
    add_text(slide, body, left + 0.20, top + 0.58, width - 0.4, height - 0.72, size=15, color=MUTED)


def add_deliverable_box(slide, left, top, width, text):
    add_rect(slide, left, top, width, 0.62, SOFT_TEAL_BG, TEAL, radius=True)
    add_text(slide, text, left + 0.15, top + 0.11, width - 0.3, 0.42, size=13, bold=True, color=NAVY, valign=MSO_ANCHOR.MIDDLE)


def add_takeaway(slide, text):
    add_rect(slide, 1.00, 5.92, 11.30, 0.62, NAVY, NAVY, radius=True)
    add_rect(slide, 1.10, 6.02, 1.15, 0.42, GOLD, GOLD, radius=True)
    add_text(
        slide,
        "Key Takeaway",
        1.10,
        6.04,
        1.15,
        0.38,
        size=10,
        bold=True,
        color=RGBColor(255, 255, 255),
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        text,
        2.38,
        6.02,
        9.70,
        0.42,
        size=13,
        color=RGBColor(255, 255, 255),
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_frame(slide, title, subtitle, page, section_idx):
    add_rect(slide, 0, 0, 13.333, 7.5, BG, BG)
    add_rect(slide, 0, 0, 13.333, 0.74, NAVY, NAVY)
    add_rect(slide, 0.48, 0.74, 0.08, 6.12, TEAL, TEAL)
    add_text(slide, title, 0.75, 0.13, 8.3, 0.30, size=24, bold=True, color=RGBColor(255, 255, 255), font="Noto Serif CJK SC")
    add_text(slide, subtitle, 9.0, 0.16, 3.25, 0.24, size=11, color=RGBColor(226, 231, 237), align=PP_ALIGN.RIGHT)

    add_rect(slide, 0.55, 6.97, 12.20, 0.02, LINE, LINE)
    add_text(slide, "詹绍基 | 孙仁组科研助理面试", 0.75, 7.00, 4.2, 0.20, size=10, color=MUTED)
    add_text(slide, f"{page:02d}", 11.95, 6.98, 0.5, 0.22, size=11, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)

    nav_left = 3.95
    nav_width = 1.85
    for idx, label in enumerate(SECTIONS):
        c = TEAL if idx == section_idx else RGBColor(219, 224, 231)
        tc = RGBColor(255, 255, 255) if idx == section_idx else MUTED
        tab = add_rect(slide, nav_left + idx * (nav_width + 0.08), 7.00, nav_width, 0.18, c, c, radius=True)
        if idx == section_idx:
            tab.line.color.rgb = GOLD
            tab.line.width = Pt(1.3)
        add_text(
            slide,
            label,
            nav_left + idx * (nav_width + 0.08),
            7.00,
            nav_width,
            0.18,
            size=9,
            bold=(idx == section_idx),
            color=tc,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )


def content_slide(prs, title, subtitle, page, sec_idx):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_frame(s, title, subtitle, page, sec_idx)
    return s


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Page 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, 13.333, 7.5, BG, BG)
    add_rect(s, 0, 0, 13.333, 1.02, NAVY, NAVY)
    add_text(s, "孙仁教授课题组科研助理岗位应聘汇报", 0.85, 0.20, 11.8, 0.50, size=34, bold=True, color=RGBColor(255, 255, 255), font="Noto Serif CJK SC")

    add_text(s, "姓名：詹绍基 | 学历：苏州大学医学影像学本科", 0.95, 1.50, 8.8, 0.36, size=22, bold=True, color=NAVY)
    add_text(s, "汇报时长：20分钟 | 18550954831 | z1309457869@163.com", 0.95, 2.05, 10.9, 0.30, size=16, color=MUTED)

    add_card(s, "今天重点", "我和岗位的匹配度\n我能给团队的价值\n入职后的执行计划", 0.95, 2.70, 3.65, 1.75)
    add_card(s, "面试定位", "围绕岗位职责回答\n不跑偏、不堆砌技术细节", 4.90, 2.70, 3.65, 1.75)
    add_card(s, "目标", "快速上手\n稳定交付\n长期深耕", 8.85, 2.70, 3.65, 1.75)

    add_takeaway(s, "本次汇报以岗位需求为中心：证明我能快速上手并持续为课题组创造价值。")

    # Page 2: Background + capability
    s = content_slide(prs, "2. 个人核心背景与能力总览", "Background and fit", 2, 0)
    add_rect(s, 1.00, 1.20, 11.30, 1.65, CARD, LINE, radius=True)
    add_text(s, "教育背景时间线", 1.25, 1.48, 2.2, 0.24, size=18, bold=True)
    add_text(s, "苏州大学 医学影像学（辅修计算机） 2019.09-2024.06 | 1年三甲医院临床实习经验", 1.25, 1.86, 10.7, 0.24, size=15)
    add_text(s, "北京师范大学 法学（第二学士学位） 2024.09-2026.07（预计）", 1.25, 2.18, 10.7, 0.24, size=15)

    add_card(s, "医学基础能力", "熟悉临床科研规范", 1.00, 3.15, 5.35, 1.95)
    add_text(s, "病毒、免疫、疫苗研发医学逻辑", 1.25, 3.72, 5.0, 0.24, size=15, bold=True, color=NAVY)
    add_card(s, "数据与编程能力", "支持系统生物学数据处理", 6.95, 3.15, 5.35, 1.95)
    add_text(s, "Python 统计分析与可视化", 7.20, 3.72, 5.0, 0.24, size=15, bold=True, color=NAVY)
    add_card(s, "全流程科研能力", "可闭环协同", 1.00, 5.25, 5.35, 0.95)
    add_text(s, "实验执行 - 数据处理 - 报告写作闭环", 1.25, 5.82, 5.0, 0.24, size=14, bold=True, color=NAVY)
    add_card(s, "英语与文献能力", "CET-6 595分\n可读写英文文献并支持论文协作", 6.95, 5.25, 5.35, 0.95)
    add_takeaway(s, "医学基础+数据能力+科研素养与严谨性，完全匹配课题组科研助理岗位的核心需求。")

    # Page 3: Understanding lab directions
    s = content_slide(prs, "3. 对课题组研究方向的理解与契合点", "Direction mapping", 3, 0)
    add_card(
        s,
        "方向A：病毒基因组精准功能定位与疫苗",
        "团队方向：单核苷酸精度功能图谱与精准减毒疫苗设计（Du et al., Science 2018）\n我的契合：医学理解 + 数据分析，支撑筛选数据统计与解读",
        1.00,
        1.35,
        3.55,
        4.75,
    )
    add_card(
        s,
        "方向B：感染后与接种后免疫反应评价",
        "团队方向：抗体特异性与T/B细胞机制系统评估\n我的契合：临床免疫检测逻辑 + 结果可视化能力",
        4.88,
        1.35,
        3.55,
        4.75,
    )
    add_card(
        s,
        "方向C：蛋白结构与功能研究",
        "团队方向：体外进化+结构解析+质谱，开发蛋白/多肽结合体（Dai X et al., Nature 2017/2018）\n我的契合：严谨对照思维 + 数据质控与分析协同",
        8.76,
        1.35,
        3.55,
        4.75,
    )
    add_takeaway(s, "我的能力体系可全面适配课题组三大方向，可快速融入并支撑课题推进。")

    # Page 4: transferability proof
    s = content_slide(prs, "4. 核心科研实践与能力证明", "Transferable capability", 4, 1)
    add_rect(s, 1.00, 1.20, 11.30, 4.95, CARD, LINE, radius=True)
    add_text(s, "实践一：科研全流程执行能力", 1.25, 1.55, 4.5, 0.24, size=18, bold=True)
    add_bullets(
        s,
        [
            "独立完成实验设计、对照实验、数据统计、图表与报告输出。",
            "基于Python完成全流程数据统计、多维度分析与学术图表可视化。",
            "产出4份完整实验报告、10+组学术图表，过程可复查可复现。",
            "可迁移到岗位：实验执行、记录归档、数据分析与结果汇报。",
        ],
        1.25,
        1.90,
        10.7,
        1.45,
        size=15,
    )
    add_text(s, "实践二：医院临床实习与科研辅助", 1.25, 3.70, 4.8, 0.24, size=18, bold=True)
    add_bullets(
        s,
        [
            "1年三甲医院实习，熟悉临床科研规范、伦理与合规要求。",
            "参与数据整理、统计分析和科研图表绘制。",
            "可迁移到岗位：规范实验流程、严谨记录、合规协作。",
        ],
        1.25,
        4.03,
        10.7,
        1.45,
        size=15,
    )
    add_takeaway(s, "我已具备科研助理所需的实验执行、数据处理、报告撰写、规范管理全流程能力，入职即可上手。")

    # Page 5: requirements mapping
    s = content_slide(prs, "5. 岗位职责与应聘要求100%对标", "Requirement checklist", 5, 1)
    add_card(s, "职责1：参与生物实验", "可通过系统学习快速掌握分子克隆、流式、小鼠实验等核心技术", 1.00, 1.35, 5.35, 1.95)
    add_text(s, "实验执行与规范记录能力", 1.25, 1.92, 5.0, 0.24, size=15, bold=True, color=NAVY)
    add_card(s, "职责2：记录+统计+分析", "可按时提交记录报告，并形成稳定数据输出", 6.95, 1.35, 5.35, 1.95)
    add_text(s, "独立完成数据统计分析", 7.20, 1.92, 5.0, 0.24, size=15, bold=True, color=NAVY)
    add_card(s, "职责3：论文与综述", "可协助文献检索、图表绘制、结果段落整理", 1.00, 3.55, 5.35, 1.95)
    add_text(s, "协助论文写作流程", 1.25, 4.12, 5.0, 0.24, size=15, bold=True, color=NAVY)
    add_card(s, "职责4：PI交办协作任务", "匹配：主动负责、沟通顺畅、按节点交付\n可稳定支持团队协同", 6.95, 3.55, 5.35, 1.95)
    add_text(s, "硬条件核对：本科及以上学历 / 英文文献能力 / 责任心与团队协作，均满足。", 1.10, 5.72, 11.1, 0.22, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_takeaway(s, "岗位职责和应聘条件逐条对齐，我可以全面胜任科研助理的各项核心工作。")

    # Page 6: 0-6 month plan
    s = content_slide(prs, "6. 入职后0-6个月工作规划", "Execution roadmap", 6, 2)
    add_card(s, "第1个月：快速融入期", "熟悉实验室SOP、设备与规章\n精读代表性论文5-10篇\n完成基础辅助工作", 1.00, 1.35, 3.55, 4.75)
    add_card(s, "第2-3个月：技能上手期", "掌握基础生物实验技术\n完成数据整理、统计分析与图表\n协助文献梳理", 4.88, 1.35, 3.55, 4.75)
    add_card(s, "第4-6个月：独立执行期", "承担子实验模块并规范交付\n协助论文图表与初稿材料\n持续输出文献综述", 8.76, 1.35, 3.55, 4.75)
    add_deliverable_box(s, 1.22, 5.05, 3.10, "交付物：SOP学习笔记 + 代表论文精读报告")
    add_deliverable_box(s, 5.10, 5.05, 3.10, "交付物：规范实验记录模板 + 实验数据分析报告")
    add_deliverable_box(s, 8.98, 5.05, 3.10, "交付物：独立子实验完整报告 + 领域综述初稿")
    add_takeaway(s, "执行路径清晰可落地：快速融入、稳步上手、形成独立可交付能力。")

    # Page 7: long-term goal
    s = content_slide(prs, "7. 长期发展目标", "Team-first growth", 7, 2)
    add_rect(s, 1.00, 1.20, 11.30, 4.95, CARD, LINE, radius=True)
    add_bullets(
        s,
        [
            "长期目标：深耕病毒学、疫苗研发与系统免疫组学方向，未来攻读本方向博士。",
            "优先路径：先在科研助理阶段完成高质量交付，优先申请孙仁老师团队博士，基于科研助理阶段的成果积累实现无缝衔接。",
            "工作原则：团队任务第一优先，用实际成果证明价值，再推进个人成长。",
            "底线承诺：无论后续发展，均完成手头任务并做好项目交接、数据和文档沉淀。",
        ],
        1.20,
        1.65,
        10.9,
        3.8,
        size=16,
    )
    add_takeaway(s, "以团队发展为第一优先，规划清晰且稳定，可与课题组实现双赢成长。")

    # Page 8: thanks
    s = content_slide(prs, "8. 致谢与Q&A", "Thanks", 8, 3)
    add_rect(s, 1.00, 1.55, 11.30, 4.65, CARD, LINE, radius=True)
    add_text(s, "感谢孙老师和各位老师的聆听！欢迎批评指正", 1.25, 2.05, 10.8, 0.65, size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font="Noto Serif CJK SC")
    add_text(s, "詹绍基 | 18550954831 | z1309457869@163.com", 1.25, 3.20, 10.8, 0.30, size=20, color=TEXT, align=PP_ALIGN.CENTER)
    add_text(s, "可提供过往实验的完整报告、代码与原始数据，确保全流程可复现。", 1.25, 4.08, 10.8, 0.25, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    add_takeaway(s, "期待加入团队，全力配合科研工作，在稳定交付中持续成长。")

    # Backup A
    s = content_slide(prs, "备用页A：过往科研成果详细数据与报告摘要", "Backup A", 9, 3)
    add_rect(s, 1.00, 1.20, 11.30, 5.30, CARD, LINE, radius=True)
    add_bullets(
        s,
        [
            "实验报告摘要：4份完整实验报告（问题定义、实验设计、结果分析、结论）。",
            "数据证据：10+组学术图表与统计表格，支持结果追溯。",
            "流程证据：实验记录模板、数据清洗脚本、结果复现说明。",
            "追问用途：证明科研能力、数据处理能力、实验严谨性。",
        ],
        1.25,
        1.70,
        10.8,
        4.6,
        size=15,
    )
    add_takeaway(s, "备用页A用于快速回应“你能不能独立做事、做得是否严谨”的追问。")

    # Backup B
    s = content_slide(prs, "备用页B：领域核心文献阅读清单", "Backup B", 10, 3)
    add_rect(s, 1.00, 1.20, 11.30, 5.30, CARD, LINE, radius=True)
    add_bullets(
        s,
        [
            "孙仁团队代表作：Science 2018（流感疫苗设计）；Nature 2017/2018（结构解析）。",
            "病毒功能图谱：mBio 2017；PNAS 2017；Cell Host Microbe 2016；NPJ Vaccines 2020。",
            "阅读脉络：病毒基因组功能定位 -> 免疫逃逸机制 -> 疫苗设计策略。",
            "追问用途：证明我对团队研究脉络做过系统准备，不是临时应付。",
        ],
        1.25,
        1.70,
        10.8,
        4.6,
        size=15,
    )
    add_takeaway(s, "备用页B用于回答“对团队方向理解深度与文献功课”的追问。")

    prs.save(str(OUT))
    print(f"saved: {OUT}")


if __name__ == "__main__":
    build()
